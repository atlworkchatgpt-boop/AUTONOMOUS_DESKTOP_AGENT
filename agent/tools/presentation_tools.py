
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt


def create_powerpoint(
    output_path,
    title,
    slides,
    pictures=None
):
    """
    Create a PowerPoint presentation.

    slides:
        [
            {
                "title": "Slide title",
                "body": "Slide text"
            }
        ]

    pictures:
        [
            {
                "slide": 1,
                "path": "C:/path/photo.jpg"
            }
        ]
    """

    output = Path(output_path)
    output.parent.mkdir(parents=True, exist_ok=True)

    prs = Presentation()

    # Title slide
    layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(layout)

    slide.shapes.title.text = title

    if slide.placeholders:
        for placeholder in slide.placeholders:
            if placeholder.placeholder_format.idx != 0:
                try:
                    placeholder.text = "Created by Autonomous Desktop AI"
                    break
                except Exception:
                    pass

    # Content slides
    for item in slides or []:
        layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(layout)

        slide.shapes.title.text = str(
            item.get("title", "Untitled")
        )

        body = item.get("body", "")

        if len(slide.placeholders) > 1:
            frame = slide.placeholders[1].text_frame
            frame.clear()

            for line in str(body).splitlines():
                p = frame.add_paragraph()
                p.text = line
                p.font.size = Pt(20)

    # Pictures
    for picture in pictures or []:
        try:
            slide_number = int(picture.get("slide", 1))
            path = Path(picture["path"])

            if not path.exists():
                continue

            if slide_number < 1 or slide_number > len(prs.slides):
                continue

            slide = prs.slides[slide_number - 1]

            slide.shapes.add_picture(
                str(path),
                Inches(5.4),
                Inches(1.6),
                width=Inches(3.3)
            )

        except Exception as exc:
            print("[PPT IMAGE ERROR]", exc)

    prs.save(str(output))

    return str(output)


if __name__ == "__main__":
    result = create_powerpoint(
        "ADA_created_presentation.pptx",
        "Autonomous Desktop AI",
        [
            {
                "title": "ADA",
                "body": "Autonomous Desktop AI"
            },
            {
                "title": "Capabilities",
                "body": "Computer automation\nFile operations\nPowerPoint creation\nAI planning"
            }
        ]
    )

    print("Created:", result)
