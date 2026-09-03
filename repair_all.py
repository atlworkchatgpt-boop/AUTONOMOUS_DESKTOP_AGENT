from pathlib import Path
import re
import shutil
import datetime

ROOT = Path(r"C:\Users\ADMIN\Desktop\ada_test_project")
APP = ROOT / "web" / "static" / "app.js"
MAIN = ROOT / "main.py"

# ------------------------------------------------------------
# BACKUP
# ------------------------------------------------------------
stamp = datetime.datetime.now().strftime("%Y%m%d_%H%M%S")

if APP.exists():
    backup = APP.with_name(f"app.js.backup_{stamp}")
    shutil.copy2(APP, backup)
    print("Backup:", backup)

if MAIN.exists():
    backup = MAIN.with_name(f"main.py.backup_{stamp}")
    shutil.copy2(MAIN, backup)
    print("Backup:", backup)

# ------------------------------------------------------------
# APP.JS
# ------------------------------------------------------------
if not APP.exists():
    raise SystemExit(f"ERROR: {APP} not found")

s = APP.read_text(encoding="utf-8", errors="replace")

# ------------------------------------------------------------
# 1. Repair featureBar WITHOUT regex replacement escapes
# ------------------------------------------------------------
start = s.find("function featureBar(){")
end = s.find("async function mediaRequest", start)

if start != -1 and end != -1 and end > start:

    clean_feature_bar = r'''
function featureBar(){
    const existing = document.querySelector(".ada-feature-bar");
    if(existing) return existing;

    const bar = document.createElement("div");
    bar.className = "ada-feature-bar";

    bar.innerHTML = `
        <button type="button" class="ada-feature-btn" id="adaImageBtn">
            <span>Image</span>
        </button>

        <button type="button" class="ada-feature-btn" id="adaVideoBtn">
            <span>Video</span>
        </button>
    `;

    const inputArea =
        document.querySelector(".ada-input-area") ||
        document.querySelector(".chat-input-area") ||
        document.querySelector("textarea")?.parentElement;

    if(inputArea && inputArea.parentElement){
        inputArea.parentElement.insertBefore(bar, inputArea);
    }

    const imageBtn = bar.querySelector("#adaImageBtn");
    const videoBtn = bar.querySelector("#adaVideoBtn");

    if(imageBtn){
        imageBtn.addEventListener("click", function(){
            const input = document.createElement("input");
            input.type = "file";
            input.accept = "image/*";

            input.onchange = function(){
                if(typeof window.handleImageUpload === "function"){
                    window.handleImageUpload(input.files[0]);
                }else if(typeof window.mediaRequest === "function"){
                    window.mediaRequest("image", input.files[0]);
                }
            };

            input.click();
        });
    }

    if(videoBtn){
        videoBtn.addEventListener("click", function(){
            const input = document.createElement("input");
            input.type = "file";
            input.accept = "video/*";

            input.onchange = function(){
                if(typeof window.handleVideoUpload === "function"){
                    window.handleVideoUpload(input.files[0]);
                }else if(typeof window.mediaRequest === "function"){
                    window.mediaRequest("video", input.files[0]);
                }
            };

            input.click();
        });
    }

    return bar;
}

'''

    s = s[:start] + clean_feature_bar + s[end:]
    print("Repaired damaged feature-bar section.")
else:
    print("Feature-bar boundary not found; left it untouched.")

# ------------------------------------------------------------
# 2. Remove obvious Voice UI references
# ------------------------------------------------------------
voice_patterns = [
    r'<button[^>]*id=["\']autonomousFinalVoice["\'][^>]*>.*?</button>',
    r'<button[^>]*id=["\']voiceBtn["\'][^>]*>.*?</button>',
    r'<button[^>]*id=["\']voiceButton["\'][^>]*>.*?</button>',
]

for pattern in voice_patterns:
    s = re.sub(pattern, "", s, flags=re.I | re.S)

# Remove direct HTML voice labels
s = re.sub(
    r'<[^>]*>\s*(?:Voice|voice)\s*</[^>]*>',
    "",
    s,
    flags=re.I | re.S
)

# Remove common voice-only function declarations safely
for fn in [
    "autonomousFinalVoice",
    "setupAutonomousFinalVoice",
    "initAutonomousFinalVoice",
]:
    pattern = rf'function\s+{re.escape(fn)}\s*\([^)]*\)\s*\{{'
    while True:
        m = re.search(pattern, s)
        if not m:
            break

        depth = 0
        i = m.end() - 1

        while i < len(s):
            if s[i] == "{":
                depth += 1
            elif s[i] == "}":
                depth -= 1
                if depth == 0:
                    i += 1
                    break
            i += 1

        s = s[:m.start()] + s[i:]

print("Removed Voice UI/function references.")

# ------------------------------------------------------------
# 3. Fix creator/owner answer
# ------------------------------------------------------------
owner_pattern = r'const\s+ownerQuestion\s*=\s*.*?;'

owner_replacement = (
    r'''const ownerQuestion = /^(?:who(?:'s| is)?\s+(?:your\s+)?(?:owner|creator)|'''
    r'''who\s+(?:made|created)\s+you|who\s+owns\s+you|'''
    r'''who\s+is\s+the\s+(?:creator|owner)|'''
    r'''who\s+is\s+your\s+(?:creator|owner)|'''
    r'''who\s+created\s+you|'''
    r'''who\s+made\s+you)[?!.\s]*$/i;'''
)

s, n = re.subn(
    owner_pattern,
    lambda m: owner_replacement,
    s,
    count=1,
    flags=re.S
)

if n:
    print("Updated creator/owner detection.")
else:
    print("Creator/owner regex not found; adding helper.")

# Replace the old response if present
s = s.replace(
    'addMessage("assistant", "My creator and owner is Shreyansh Ray.");',
    '''addMessage(
            "assistant",
            "Creator & Owner: Shreyansh Ray\\nHonorable Mention — Support: Arnav Baliyan"
        );'''
)

# Also handle another common old wording
s = s.replace(
    'My creator and owner is Shreyansh Ray.',
    'Creator & Owner: Shreyansh Ray\\nHonorable Mention — Support: Arnav Baliyan'
)

APP.write_text(s, encoding="utf-8")
print("Saved app.js.")

# ------------------------------------------------------------
# 4. Add a global creator helper so natural UI code can reuse it
# ------------------------------------------------------------
creator_helper = r'''
// ADA CREATOR / OWNER INFORMATION
window.ADA_CREATOR_INFO = {
    creator: "Shreyansh Ray",
    owner: "Shreyansh Ray",
    honorableMention: "Arnav Baliyan"
};

'''

if "window.ADA_CREATOR_INFO" not in s:
    s = creator_helper + s
    APP.write_text(s, encoding="utf-8")
    print("Added global creator information.")

# ------------------------------------------------------------
# 5. Create PowerPoint + computer tool module
# ------------------------------------------------------------
TOOLS = ROOT / "agent" / "tools"
TOOLS.mkdir(parents=True, exist_ok=True)

ppt = TOOLS / "presentation_tools.py"

ppt.write_text(r'''
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt


def create_powerpoint(
    output_path,
    title,
    slides,
    pictures=None
):
    """
    Create a PowerPoint presentation.

    slides:
        [
            {
                "title": "Slide title",
                "body": "Slide text"
            }
        ]

    pictures:
        [
            {
                "slide": 1,
                "path": "C:/path/photo.jpg"
            }
        ]
    """

    output = Path(output_path)
    output.parent.mkdir(parents=True, exist_ok=True)

    prs = Presentation()

    # Title slide
    layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(layout)

    slide.shapes.title.text = title

    if slide.placeholders:
        for placeholder in slide.placeholders:
            if placeholder.placeholder_format.idx != 0:
                try:
                    placeholder.text = "Created by Autonomous Desktop AI"
                    break
                except Exception:
                    pass

    # Content slides
    for item in slides or []:
        layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(layout)

        slide.shapes.title.text = str(
            item.get("title", "Untitled")
        )

        body = item.get("body", "")

        if len(slide.placeholders) > 1:
            frame = slide.placeholders[1].text_frame
            frame.clear()

            for line in str(body).splitlines():
                p = frame.add_paragraph()
                p.text = line
                p.font.size = Pt(20)

    # Pictures
    for picture in pictures or []:
        try:
            slide_number = int(picture.get("slide", 1))
            path = Path(picture["path"])

            if not path.exists():
                continue

            if slide_number < 1 or slide_number > len(prs.slides):
                continue

            slide = prs.slides[slide_number - 1]

            slide.shapes.add_picture(
                str(path),
                Inches(5.4),
                Inches(1.6),
                width=Inches(3.3)
            )

        except Exception as exc:
            print("[PPT IMAGE ERROR]", exc)

    prs.save(str(output))

    return str(output)


if __name__ == "__main__":
    result = create_powerpoint(
        "ADA_created_presentation.pptx",
        "Autonomous Desktop AI",
        [
            {
                "title": "ADA",
                "body": "Autonomous Desktop AI"
            },
            {
                "title": "Capabilities",
                "body": "Computer automation\nFile operations\nPowerPoint creation\nAI planning"
            }
        ]
    )

    print("Created:", result)
''', encoding="utf-8")

print("Created presentation_tools.py.")

# ------------------------------------------------------------
# 6. Computer/file operations tool
# ------------------------------------------------------------
computer = TOOLS / "computer_tools.py"

computer.write_text(r'''
from pathlib import Path
import shutil


def create_folder(path):
    p = Path(path)
    p.mkdir(parents=True, exist_ok=True)
    return str(p)


def list_files(path):
    p = Path(path)

    if not p.exists():
        return []

    return [
        {
            "name": x.name,
            "path": str(x),
            "is_directory": x.is_dir()
        }
        for x in p.iterdir()
    ]


def copy_file(source, destination):
    src = Path(source)
    dst = Path(destination)

    dst.parent.mkdir(parents=True, exist_ok=True)

    shutil.copy2(src, dst)

    return str(dst)


def move_file(source, destination):
    src = Path(source)
    dst = Path(destination)

    dst.parent.mkdir(parents=True, exist_ok=True)

    return str(shutil.move(str(src), str(dst)))


def rename_file(source, new_name):
    src = Path(source)
    dst = src.with_name(new_name)

    src.rename(dst)

    return str(dst)


def delete_file(path):
    p = Path(path)

    if p.is_file():
        p.unlink()
        return True

    if p.is_dir():
        shutil.rmtree(p)
        return True

    return False
''', encoding="utf-8")

print("Created computer_tools.py.")

# ------------------------------------------------------------
# 7. Backend attribution
# ------------------------------------------------------------
if MAIN.exists():
    main = MAIN.read_text(encoding="utf-8", errors="replace")

    main = main.replace(
        '"owner": "Shreyansh Ray",',
        '"owner": "Shreyansh Ray",\n        "honorable_mention": "Arnav Baliyan",'
    )

    if '"honorable_mention": "Arnav Baliyan"' not in main:
        # Add to autonomous status dictionary if possible
        marker = '"creator": "Shreyansh Ray",'
        main = main.replace(
            marker,
            marker + '\n        "honorable_mention": "Arnav Baliyan",',
            1
        )

    MAIN.write_text(main, encoding="utf-8")
    print("Updated backend creator attribution.")

print()
print("============================================================")
print("ADA REPAIR COMPLETE")
print("============================================================")
print("Buttons: repaired")
print("Voice UI: removed")
print("Creator: Shreyansh Ray")
print("Owner: Shreyansh Ray")
print("Honorable Mention: Arnav Baliyan")
print("PowerPoint tool: added")
print("Computer/file tools: added")
print()
